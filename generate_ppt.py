from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color Palette ───────────────────────────────────────────────
GREEN_DARK   = RGBColor(0x16, 0xa3, 0x4a)   # #16a34a
GREEN_LIGHT  = RGBColor(0xdc, 0xfc, 0xe7)   # #dcfce7
GREEN_MID    = RGBColor(0x22, 0xc5, 0x5e)   # #22c55e
BLUE_DARK    = RGBColor(0x1d, 0x4e, 0xd8)   # #1d4ed8
BLUE_MID     = RGBColor(0x3b, 0x82, 0xf6)   # #3b82f6
BLUE_LIGHT   = RGBColor(0xbe, 0xfb, 0xff)   # #bfdbfe (light blue)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x0f, 0x17, 0x2a)
LIGHT_TEXT   = RGBColor(0xf0, 0xfd, 0xf4)
ACCENT       = RGBColor(0x06, 0xb6, 0xd4)   # cyan

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helper functions ─────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_w:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_textbox_multiline(slide, lines, l, t, w, h,
                          font_size=16, bold=False, color=DARK_TEXT,
                          align=PP_ALIGN.LEFT, line_spacing=1.0):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def gradient_bg(slide, top_color, bot_color):
    """Approximate gradient with two overlapping rects"""
    add_rect(slide, 0, 0, 13.33, 3.75, fill_rgb=top_color)
    add_rect(slide, 0, 3.75, 13.33, 3.75, fill_rgb=bot_color)

def add_image(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def section_header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=GREEN_DARK)
    add_text(slide, title, 0.3, 0.1, 12.5, 0.85,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.75, 12.5, 0.4,
                 font_size=13, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# Full background split: top green, bottom light blue
add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x0f, 0x2a, 0x40))  # dark navy bg

# Left colored bar
add_rect(s1, 0, 0, 0.35, 7.5, fill_rgb=GREEN_MID)

# Top decorative strip
add_rect(s1, 0.35, 0, 12.98, 0.12, fill_rgb=ACCENT)

# Big glowing card
add_rect(s1, 0.55, 0.4, 8.1, 6.7,
         fill_rgb=RGBColor(0x05, 0x1a, 0x30),
         line_rgb=GREEN_MID, line_w=1.5)

# Event badge
add_rect(s1, 0.75, 0.6, 4.5, 0.45, fill_rgb=GREEN_DARK)
add_text(s1, "🌿  CLAUDE BASED HACKATHON  •  CUTM 2026",
         0.78, 0.62, 4.4, 0.38,
         font_size=10, bold=True, color=WHITE)

# Title
add_text(s1, "AI-Assisted Precision Nutrient\n& Water Management\nin Agriculture",
         0.65, 1.2, 7.8, 3.2,
         font_size=32, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)

# Subtitle / tagline
add_text(s1, "Powered by Machine Learning & Gemini AI  |  FarmGuru Platform",
         0.65, 4.3, 7.8, 0.55,
         font_size=13, italic=True, color=ACCENT, align=PP_ALIGN.LEFT)

# Divider line
add_rect(s1, 0.65, 4.95, 7.6, 0.04, fill_rgb=GREEN_MID)

# Members section label
add_text(s1, "Prepared by:", 0.65, 5.1, 3, 0.4,
         font_size=13, bold=True, color=GREEN_LIGHT)

members = [
    "V. Harish       (240101120077)",
    "T. Himagiri     (240101120107)",
    "B. Hemanth    (240101120091)",
    "M. Jayanth     (240101120108)",
]
for i, m in enumerate(members):
    add_text(s1, f"  ◆  {m}", 0.65, 5.52 + i * 0.38, 7.8, 0.36,
             font_size=12.5, color=WHITE)

# Right panel – hackathon details
add_rect(s1, 8.9, 0.4, 4.1, 6.7,
         fill_rgb=RGBColor(0x0a, 0x25, 0x42),
         line_rgb=BLUE_MID, line_w=1)

add_text(s1, "🏆  Event Details", 9.05, 0.6, 3.8, 0.5,
         font_size=14, bold=True, color=BLUE_MID)
add_rect(s1, 9.0, 1.1, 3.8, 0.04, fill_rgb=BLUE_MID)

event_lines = [
    "📅  Aug 18th & 19th, 2026",
    "🏫  Centurion University",
    "      of Technology &",
    "      Management, Odisha",
    "",
    "🤝  Organized by:",
    "      IEEE  •  IIC",
    "",
    "🎯  Theme:",
    "      Code · Innovate · Impact",
    "",
    "⚙️   Tech Stack:",
    "      React  •  FastAPI",
    "      Gemini AI  •  Supabase",
    "      scikit-learn  •  Python",
]
for i, ln in enumerate(event_lines):
    add_text(s1, ln, 9.05, 1.2 + i * 0.34, 3.8, 0.33,
             font_size=11, color=WHITE if ln.strip() else WHITE)

# Bottom bar
add_rect(s1, 0, 7.2, 13.33, 0.3, fill_rgb=GREEN_DARK)
add_text(s1, "FarmGuru  •  Smart Agriculture Platform  •  2026",
         0, 7.22, 13.33, 0.26,
         font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 – SRS (Software Requirements Specification)
# ════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xf4, 0xf9, 0xf4))  # light green bg

section_header(s2, "Software Requirements Specification (SRS)",
               "FarmGuru — AI-Powered Soil & Water Management Platform")

# Two columns
col_w = 5.9

# ── LEFT COLUMN ──
add_rect(s2, 0.3, 1.25, col_w, 5.9, fill_rgb=WHITE,
         line_rgb=GREEN_MID, line_w=1)
add_rect(s2, 0.3, 1.25, col_w, 0.45, fill_rgb=GREEN_DARK)
add_text(s2, "  Functional Requirements", 0.3, 1.27, col_w, 0.4,
         font_size=13, bold=True, color=WHITE)

func_reqs = [
    "FR-01  Users input N, P, K & Moisture values via web form",
    "FR-02  ML model predicts recommended crop type",
    "FR-03  Soil health score (0–100%) computed by AI model",
    "FR-04  System calculates additional water requirement (mm)",
    "FR-05  Gemini 2.5 Flash generates detailed text recommendations",
    "FR-06  All analyses stored in Supabase cloud database",
    "FR-07  Dashboard shows trend graphs & average NPK charts",
    "FR-08  AI chatbot answers farming queries with suggestions",
    "FR-09  Light / Dark theme toggle available site-wide",
    "FR-10  History table displays all past soil analyses",
]
for i, r in enumerate(func_reqs):
    add_text(s2, f"  ✔  {r}", 0.35, 1.75 + i * 0.5, col_w - 0.1, 0.48,
             font_size=11, color=DARK_TEXT)

# ── RIGHT COLUMN ──
rx = 0.3 + col_w + 0.3
add_rect(s2, rx, 1.25, col_w, 2.55, fill_rgb=WHITE,
         line_rgb=BLUE_MID, line_w=1)
add_rect(s2, rx, 1.25, col_w, 0.45, fill_rgb=BLUE_DARK)
add_text(s2, "  Non-Functional Requirements", rx, 1.27, col_w, 0.4,
         font_size=13, bold=True, color=WHITE)

nfunc = [
    "NFR-01  Response time < 3 sec for ML prediction",
    "NFR-02  AI chat streaming response within 5 seconds",
    "NFR-03  System supports concurrent multi-user access",
    "NFR-04  API keys secured via .env (not in codebase)",
]
for i, r in enumerate(nfunc):
    add_text(s2, f"  ✔  {r}", rx + 0.05, 1.77 + i * 0.54, col_w - 0.1, 0.5,
             font_size=11, color=DARK_TEXT)

# ── Tech & System Requirements ──
add_rect(s2, rx, 3.95, col_w, 3.2, fill_rgb=WHITE,
         line_rgb=ACCENT, line_w=1)
add_rect(s2, rx, 3.95, col_w, 0.45, fill_rgb=RGBColor(0x06, 0x5f, 0x46))
add_text(s2, "  System & Tech Requirements", rx, 3.97, col_w, 0.4,
         font_size=13, bold=True, color=WHITE)

tech = [
    "Frontend  :  React (Vite) + Tailwind CSS + Recharts",
    "Backend   :  Python FastAPI + Uvicorn server",
    "AI Model  :  Google Gemini 2.5 Flash (streaming SSE)",
    "ML Models:  scikit-learn (crop + farm analytics .pkl)",
    "Database  :  Supabase (PostgreSQL cloud)",
    "Hosting   :  Local (port 5173 / 8000), GitHub repo",
]
for i, r in enumerate(tech):
    add_text(s2, f"  ▶  {r}", rx + 0.05, 4.45 + i * 0.42, col_w - 0.1, 0.4,
             font_size=10.5, color=DARK_TEXT)

add_rect(s2, 0, 7.2, 13.33, 0.3, fill_rgb=GREEN_DARK)
add_text(s2, "FarmGuru  •  Smart Agriculture Platform  •  2026",
         0, 7.22, 13.33, 0.26, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 – ER DIAGRAM (drawn with shapes)
# ════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xf0, 0xfd, 0xf4))
section_header(s3, "Entity-Relationship (ER) Diagram",
               "FarmGuru — Data Model & Relationships")

# Entity helper
def entity(slide, name, attrs, l, t, w=3.0):
    h = 0.42 + len(attrs) * 0.36
    add_rect(slide, l, t, w, 0.42, fill_rgb=GREEN_DARK)
    add_text(slide, name, l + 0.1, t + 0.04, w - 0.2, 0.35,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, l, t + 0.42, w, h - 0.42, fill_rgb=WHITE,
             line_rgb=GREEN_MID, line_w=1)
    for i, a in enumerate(attrs):
        clr = BLUE_DARK if a.startswith("🔑") else DARK_TEXT
        add_text(slide, f"  {a}", l + 0.05, t + 0.46 + i * 0.36, w - 0.1, 0.34,
                 font_size=10.5, color=clr)
    return l, t, w, h

# USER entity
entity(s3, "USER",
       ["🔑 user_id (PK)", "name", "email", "role (farmer/admin)"],
       0.4, 1.4)

# ANALYSIS entity (centre)
entity(s3, "SOIL_ANALYSIS",
       ["🔑 id (PK)", "user_id (FK)", "nitrogen (mg/kg)",
        "phosphorus (mg/kg)", "potassium (mg/kg)",
        "moisture (%)", "recommended_crop",
        "health_score", "water_needed (mm)",
        "condition", "created_at"],
       4.55, 1.1, w=4.0)

# RECOMMENDATION entity
entity(s3, "AI_RECOMMENDATION",
       ["🔑 rec_id (PK)", "analysis_id (FK)", "content (text)", "created_at"],
       9.3, 1.4, w=3.6)

# CHAT_LOG entity
entity(s3, "CHAT_LOG",
       ["🔑 chat_id (PK)", "user_id (FK)", "message", "response", "timestamp"],
       0.4, 4.5, w=3.6)

# DASHBOARD entity
entity(s3, "DASHBOARD_STATS",
       ["🔑 stat_id (PK)", "user_id (FK)", "total_analyses",
        "avg_health", "avg_npk_n", "avg_npk_p", "avg_npk_k"],
       9.3, 4.5, w=3.6)

# Relationship labels (simple lines as text arrows)
def rel_label(slide, text, l, t):
    add_rect(slide, l, t, 1.6, 0.32, fill_rgb=BLUE_DARK)
    add_text(slide, text, l + 0.05, t + 0.04, 1.5, 0.26,
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rel_label(s3, "1 : N  submits", 3.1, 2.55)
rel_label(s3, "1 : N  generates", 8.6, 2.55)
rel_label(s3, "1 : N  chats", 3.1, 5.2)
rel_label(s3, "1 : 1  aggregates", 8.0, 5.2)

# Legend
add_rect(s3, 0.4, 6.6, 12.5, 0.52, fill_rgb=GREEN_LIGHT, line_rgb=GREEN_MID, line_w=0.5)
add_text(s3, "  🔑 = Primary Key   |   FK = Foreign Key   |   Supabase PostgreSQL stores all entities   |   Analysis is the central entity",
         0.5, 6.65, 12.3, 0.4, font_size=10.5, color=GREEN_DARK, align=PP_ALIGN.CENTER)

add_rect(s3, 0, 7.2, 13.33, 0.3, fill_rgb=GREEN_DARK)
add_text(s3, "FarmGuru  •  Smart Agriculture Platform  •  2026",
         0, 7.22, 13.33, 0.26, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 – FRONTEND UI
# ════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xf0, 0xfd, 0xf4))
section_header(s4, "Frontend — User Interface (React + Vite + Tailwind CSS)",
               "Live FarmGuru Web Application  |  Light Theme  |  Full Responsive Design")

ui_img = r"C:\Users\HP\.gemini\antigravity\brain\17d6a25d-f043-4b44-9e1a-673030fe74ca\.user_uploaded\media_1787132114386.png"
add_image(s4, ui_img, 0.3, 1.2, 8.5, 5.75)

# Feature callout cards
feats = [
    ("🌿", "Analysis Form", "4 inputs: N, P, K,\nMoisture with\ncolored unit badges"),
    ("📊", "Dashboard", "Live Recharts graphs\nfrom Supabase\n(trend + NPK avg)"),
    ("🤖", "AI Chatbot", "Gemini streaming\nresponses + clickable\nsuggestion pills"),
    ("🌓", "Themes", "Light / Dark mode\nwith one-click\ntoggle in navbar"),
]
for i, (icon, title, desc) in enumerate(feats):
    fy = 1.2 + i * 1.45
    add_rect(s4, 9.05, fy, 3.95, 1.35, fill_rgb=WHITE,
             line_rgb=GREEN_MID if i % 2 == 0 else BLUE_MID, line_w=1.2)
    add_text(s4, f"{icon}  {title}", 9.15, fy + 0.08, 3.7, 0.42,
             font_size=13, bold=True,
             color=GREEN_DARK if i % 2 == 0 else BLUE_DARK)
    add_text(s4, desc, 9.15, fy + 0.5, 3.7, 0.78,
             font_size=10.5, color=DARK_TEXT)

add_rect(s4, 0, 7.2, 13.33, 0.3, fill_rgb=GREEN_DARK)
add_text(s4, "FarmGuru  •  Smart Agriculture Platform  •  2026",
         0, 7.22, 13.33, 0.26, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 – BACKEND (VS Code screenshot + explanation)
# ════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xf0, 0xf8, 0xff))
section_header(s5, "Backend — FastAPI + Python + ML Models",
               "main.py  |  prediction.py  |  Gemini AI  |  Supabase Integration")

vscode_img = r"C:\Users\HP\.gemini\antigravity\brain\17d6a25d-f043-4b44-9e1a-673030fe74ca\.user_uploaded\media_1787132206618.png"
add_image(s5, vscode_img, 0.3, 1.2, 7.8, 5.75)

# Backend components panel
components = [
    ("⚡", "FastAPI", "High-performance Python\nweb framework\n(Async + CORS ready)"),
    ("🧠", "ML Models", "crop_recommender.pkl\nfarm_analytics.pkl\n(scikit-learn 1.9)"),
    ("✨", "Gemini AI", "gemini-2.5-flash\nSSE streaming chat\n& soil advice"),
    ("🗄️", "Supabase", "PostgreSQL cloud DB\nAuto-logs analyses\nhistory & stats API"),
]
for i, (icon, title, desc) in enumerate(components):
    cy = 1.2 + i * 1.45
    clr = GREEN_DARK if i % 2 == 0 else BLUE_DARK
    brd = GREEN_MID if i % 2 == 0 else BLUE_MID
    add_rect(s5, 8.35, cy, 4.65, 1.35, fill_rgb=WHITE, line_rgb=brd, line_w=1.2)
    add_text(s5, f"{icon}  {title}", 8.48, cy + 0.08, 4.4, 0.42,
             font_size=13, bold=True, color=clr)
    add_text(s5, desc, 8.48, cy + 0.5, 4.4, 0.78,
             font_size=10.5, color=DARK_TEXT)

add_rect(s5, 0, 7.2, 13.33, 0.3, fill_rgb=GREEN_DARK)
add_text(s5, "FarmGuru  •  Smart Agriculture Platform  •  2026",
         0, 7.22, 13.33, 0.26, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 – CONCLUSION & THANK YOU
# ════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x0a, 0x1a, 0x30))  # dark navy

# Animated top bar (visual only)
add_rect(s6, 0, 0, 13.33, 0.12, fill_rgb=GREEN_MID)
add_rect(s6, 0, 0.12, 13.33, 0.06, fill_rgb=BLUE_MID)

# Central glowing card
add_rect(s6, 1.5, 0.8, 10.33, 5.6,
         fill_rgb=RGBColor(0x05, 0x1a, 0x30),
         line_rgb=GREEN_MID, line_w=2)

add_text(s6, "🌿  FarmGuru", 1.6, 0.95, 10.1, 0.7,
         font_size=16, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

add_text(s6, "AI-Assisted Precision Nutrient &\nWater Management in Agriculture",
         1.6, 1.55, 10.1, 1.1,
         font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s6, 3.0, 2.7, 7.33, 0.05, fill_rgb=ACCENT)

# Conclusion bullets
conclusions = [
    "✅  Successfully built an end-to-end AI + ML powered soil analysis platform",
    "✅  Integrated Gemini 2.5 Flash for real-time streaming agricultural advice",
    "✅  Supabase cloud database captures every analysis for historical tracking",
    "✅  Farmer-friendly UI with clickable AI suggestions (accessibility first)",
    "✅  Full Light / Dark theme with responsive design across all devices",
]
for i, c in enumerate(conclusions):
    add_text(s6, c, 1.8, 2.85 + i * 0.47, 9.7, 0.44,
             font_size=12, color=WHITE, align=PP_ALIGN.LEFT)

add_rect(s6, 3.0, 5.25, 7.33, 0.05, fill_rgb=BLUE_MID)

add_text(s6, "🙏  Thank You!", 1.6, 5.4, 10.1, 0.65,
         font_size=30, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

add_text(s6, "We look forward to your valuable feedback.",
         1.6, 5.95, 10.1, 0.4,
         font_size=14, italic=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

# Team credits at bottom of card
add_text(s6,
         "V. Harish (077)   ·   T. Himagiri (107)   ·   B. Hemanth (091)   ·   M. Jayanth (108)",
         1.6, 6.35, 10.1, 0.4,
         font_size=11, color=RGBColor(0x86, 0xef, 0xac), align=PP_ALIGN.CENTER)

add_rect(s6, 0, 7.2, 13.33, 0.3, fill_rgb=GREEN_DARK)
add_text(s6, "FarmGuru  •  Smart Agriculture Platform  •  2026",
         0, 7.22, 13.33, 0.26, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ─── SAVE ─────────────────────────────────────────────────────────
out = r"C:\Users\HP\.gemini\antigravity\brain\17d6a25d-f043-4b44-9e1a-673030fe74ca\FarmGuru_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
